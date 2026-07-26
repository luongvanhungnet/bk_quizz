from pathlib import Path

from docx import Document
from pptx import Presentation

from app.services.document_parser import DocumentParser


def test_docx_parser_keeps_heading_table_and_text(tmp_path: Path) -> None:
    path = tmp_path / "lesson.docx"
    document = Document()
    document.add_heading("Chương 1", level=1)
    document.add_paragraph("Nội dung bài học")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    document.save(path)

    sections = DocumentParser().parse(path)

    assert sections[0].heading == "Chương 1"
    assert "Nội dung bài học" in sections[0].text
    assert "A | B" in sections[0].text


def test_pptx_parser_keeps_slide_number_and_ignores_notes(tmp_path: Path) -> None:
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Đại số"
    slide.placeholders[1].text = "Ma trận và định thức"
    slide.notes_slide.notes_text_frame.text = "Không được index ghi chú này"
    presentation.save(path)

    sections = DocumentParser().parse(path)

    assert sections[0].heading == "Đại số"
    assert sections[0].slide_number == 1
    assert "Ma trận" in sections[0].text
    assert "ghi chú" not in sections[0].text
