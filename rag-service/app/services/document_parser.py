import re
from pathlib import Path

import pymupdf
from docx import Document as WordDocument
from pptx import Presentation

from app.models.document import DocumentSection

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}
HEADING_PATTERN = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*#*\s*$")


class DocumentParser:
    def parse(self, path: Path, original_filename: str | None = None) -> list[DocumentSection]:
        suffix = Path(original_filename or path.name).suffix.casefold()
        if suffix == ".pdf":
            sections = self._parse_pdf(path)
        elif suffix == ".docx":
            sections = self._parse_docx(path)
        elif suffix == ".pptx":
            sections = self._parse_pptx(path)
        elif suffix == ".txt":
            sections = self._parse_text(path)
        elif suffix in {".md", ".markdown"}:
            sections = self._parse_markdown(path)
        else:
            raise ValueError(f"Định dạng tài liệu không được hỗ trợ: {path.name}")
        if not any(section.text.strip() for section in sections):
            raise ValueError(f"Tài liệu {path.name} không có văn bản có thể lập chỉ mục.")
        return sections

    @staticmethod
    def _parse_docx(path: Path) -> list[DocumentSection]:
        try:
            document = WordDocument(path)
        except Exception as error:
            raise ValueError(f"Không thể đọc DOCX {path.name}.") from error
        sections: list[DocumentSection] = []
        heading: str | None = None
        lines: list[str] = []

        def flush() -> None:
            text = "\n".join(lines).strip()
            if text:
                sections.append(DocumentSection(None, heading, text))
            lines.clear()

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = (paragraph.style.name or "").casefold()
            if style_name.startswith("heading"):
                flush()
                heading = text
            else:
                lines.append(text)
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    lines.append(" | ".join(cells))
        flush()
        return sections

    @staticmethod
    def _parse_pptx(path: Path) -> list[DocumentSection]:
        try:
            presentation = Presentation(path)
        except Exception as error:
            raise ValueError(f"Không thể đọc PPTX {path.name}.") from error
        sections: list[DocumentSection] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = slide.shapes.title.text.strip() if slide.shapes.title else None
            texts: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                value = shape.text.strip()
                if value and value != title:
                    texts.append(value)
            text = "\n".join(texts).strip()
            if title and not text:
                text = title
            if text:
                sections.append(DocumentSection(None, title, text, slide_number))
        return sections

    @staticmethod
    def _parse_text(path: Path) -> list[DocumentSection]:
        return [DocumentSection(None, None, path.read_text(encoding="utf-8-sig"))]

    @staticmethod
    def _parse_pdf(path: Path) -> list[DocumentSection]:
        sections: list[DocumentSection] = []
        try:
            with pymupdf.open(path) as document:
                for index, page in enumerate(document):
                    text = page.get_text("text", sort=True).strip()
                    if text:
                        sections.append(DocumentSection(index + 1, None, text))
        except (pymupdf.FileDataError, RuntimeError) as error:
            raise ValueError(f"Không thể đọc PDF {path.name}.") from error
        return sections

    @staticmethod
    def _parse_markdown(path: Path) -> list[DocumentSection]:
        sections: list[DocumentSection] = []
        heading: str | None = None
        lines: list[str] = []
        for line in path.read_text(encoding="utf-8-sig").splitlines():
            match = HEADING_PATTERN.match(line)
            if match:
                text = "\n".join(lines).strip()
                if text:
                    sections.append(DocumentSection(None, heading, text))
                heading = match.group(1).strip()
                lines = []
            else:
                lines.append(line)
        text = "\n".join(lines).strip()
        if text:
            sections.append(DocumentSection(None, heading, text))
        return sections
